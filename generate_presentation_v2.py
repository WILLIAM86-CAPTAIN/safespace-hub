from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from PIL import Image

presentation = Presentation()

# Title slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "THE SAFE SPACE HUB"
try:
    slide.placeholders[1].text = "Project Documentation"
except Exception:
    pass

# Slide: Abstract
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Abstract"
body = slide.shapes.placeholders[1].text_frame
body.text = "THE SAFE SPACE HUB is a mental wellness web prototype focused on accessible mood tracking, assessments, and resources to support young adults."
body.add_paragraph().text = "This document presents the current system, proposed improvements, architecture, and future directions."

# Slide: Existing System
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Existing System"
body = slide.shapes.placeholders[1].text_frame
body.text = "Current demo uses frontend-only flows with localStorage for user sessions and data."
body.add_paragraph().text = "Limitations: no persistent authenticated backend, limited analytics, demo user store."

# Slide: Proposed System
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Proposed System"
body = slide.shapes.placeholders[1].text_frame
body.text = "Move to full-stack with secure SQLite-backed API, proper authentication, and data export for clinicians."
body.add_paragraph().text = "Add analytics, user roles, encrypted storage, and professional integration (telehealth links)."

# Slide: Advantages
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Advantages"
body = slide.shapes.placeholders[1].text_frame
body.text = "Benefits of the proposed system"
for bullet in [
    "Persistent, secure user data storage",
    "Scalable analytics for population insights",
    "Improved clinical integration and export",
    "Better privacy and compliance controls"
]:
    p = body.add_paragraph(); p.text = bullet; p.level = 1

# Slide: Disadvantages
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Disadvantages"
body = slide.shapes.placeholders[1].text_frame
body.text = "Potential drawbacks and trade-offs"
for bullet in [
    "Higher development and maintenance cost",
    "Requires secure hosting and compliance work",
    "Need for clinician oversight and validation",
    "Possible user privacy concerns to manage"
]:
    p = body.add_paragraph(); p.text = bullet; p.level = 1

# Slide: Software Requirements
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Software Requirements"
rows, cols = 5, 2
left = Inches(0.5); top = Inches(1.5); width = Inches(9); height = Inches(2.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.cell(0,0).text = "Component"; table.cell(0,1).text = "Requirement"
reqs = [
    ["Frontend", "HTML5, CSS3, JavaScript"],
    ["Backend", "Node.js, Express"],
    ["Database", "SQLite (or production DB)"],
    ["Runtime", "Node 14+; Python only for tooling"],
]
for i, r in enumerate(reqs, start=1):
    table.cell(i,0).text = r[0]; table.cell(i,1).text = r[1]

# Slide: About Project
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "About Project"
body = slide.shapes.placeholders[1].text_frame
body.text = "Author: Sage Williamz"
body.add_paragraph().text = "Version: 1.0.0"
body.add_paragraph().text = "Purpose: Provide a supportive, accessible prototype for young adults to track moods and access resources."

# Slide: System Architecture
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "System Architecture"
body = slide.shapes.placeholders[1].text_frame
body.text = "Client (Browser) <-> Node.js API <-> SQLite (persistent store)"
body.add_paragraph().text = "PWA-ready frontend, optional service worker, and local demo fallback."
# Add architecture image if present
image_path = Path('PROJECT IMAGE 1.webp')
if image_path.exists():
    try:
        slide.shapes.add_picture(str(image_path), Inches(6), Inches(1.5), width=Inches(3))
    except Exception:
        try:
            img = Image.open(image_path); png_path = Path('project_image_converted.png'); img.save(png_path); slide.shapes.add_picture(str(png_path), Inches(6), Inches(1.5), width=Inches(3))
        except Exception:
            pass

# Slide: Conclusion
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Conclusion"
body = slide.shapes.placeholders[1].text_frame
body.text = "THE SAFE SPACE HUB is a promising prototype that can become a clinically valuable tool with further development."
body.add_paragraph().text = "Next steps focus on security, data persistence, and clinician collaboration."

# Slide: Future Work
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Future Work"
body = slide.shapes.placeholders[1].text_frame
body.text = "Planned improvements"
for bullet in [
    "Implement secure authentication and role-based access",
    "Add encrypted database storage and backup",
    "Integrate telehealth and clinician dashboards",
    "Run validation studies and usability testing"
]:
    p = body.add_paragraph(); p.text = bullet; p.level = 1

presentation.save("THE_SAFE_SPACE_HUB_Presentation.pptx")
print("Presentation created: THE_SAFE_SPACE_HUB_Presentation.pptx")
