from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path
from PIL import Image

presentation = Presentation()

# Title slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "THE SAFE SPACE HUB"
slide.placeholders[1].text = "Mental Wellness Support Platform\nProject Overview, Features, and Insights"

# Slide 2: Project Overview
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Project Overview"
body = slide.shapes.placeholders[1].text_frame
body.text = "A mental wellness web app designed to create a safe, accessible environment for mood tracking, assessments, and supportive resources."
body.add_paragraph().text = "Built with a modern UI, secure login flow, wellness tools, and offline-friendly PWA readiness."
body.add_paragraph().text = "Designed to support users with HIPAA-style privacy, WCAG accessibility, and responsive web experience."

# Slide 3: Key Features
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Key Features"
body = slide.shapes.placeholders[1].text_frame
body.text = "Core functionality included in the Safe Space Hub"
features = [
    "Secure login and session persistence",
    "Mood tracking with user feedback",
    "Mental health assessments (PHQ-9, GAD-7)",
    "AI companion & emergency support prompts",
    "Resources, privacy policy, and accessibility focus"
]
for feature in features:
    p = body.add_paragraph()
    p.text = feature
    p.level = 1

# Slide 4: Architecture & Tech Stack
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Architecture & Tech Stack"
rows, cols = 6, 2
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(2.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
headers = ["Layer", "Technology"]
values = [
    ["Frontend", "HTML5, CSS3, JavaScript"],
    ["Backend", "Node.js + Express"],
    ["Data Storage", "localStorage demo / SQLite backend"],
    ["Design", "Responsive UI, WCAG, PWA-ready"],
    ["Deployment", "Local host launch via start-safe-space.bat"]
]
for col_idx, header in enumerate(headers):
    table.cell(0, col_idx).text = header
for row_idx, row_data in enumerate(values, start=1):
    table.cell(row_idx, 0).text = row_data[0]
    table.cell(row_idx, 1).text = row_data[1]
for row in table.rows:
    for cell in row.cells:
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)

# Slide 5: Data Insights + Charts
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Mental Health Rates for Ages 20-35"
# Bar chart data for illness rates
chart_data = ChartData()
chart_data.categories = ["PTSD", "Autism", "High Anxiety", "Bipolar"]
chart_data.add_series("Rate (%)", (12, 9, 28, 7))
x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(4.5), Inches(3)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
# Pie chart data for Uganda youth mental health vs healthy population
pie_data = ChartData()
pie_data.categories = ["Affected by Selected Illnesses", "Healthy Young Population"]
pie_data.add_series("Uganda Youth 20-35", (18, 82))
x2, y2, cx2, cy2 = Inches(5.5), Inches(1.5), Inches(4), Inches(3)
slide.shapes.add_chart(XL_CHART_TYPE.PIE, x2, y2, cx2, cy2, pie_data)

# Slide 6: UI & Experience
slide_layout = presentation.slide_layouts[6]
slide = presentation.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "User Interface & Experience"
title_frame.paragraphs[0].font.size = Pt(32)
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5), Inches(4.5))
text_frame = text_box.text_frame
text_frame.text = "The Safe Space Hub delivers an inviting, supportive interface with"
for bullet in [
    "Warm branding and thoughtful color palette",
    "Clear login and onboarding flow",
    "Accessible layout with large buttons and readable text",
    "Modal privacy policy, demo access, and emergency guidance"
]:
    p = text_frame.add_paragraph()
    p.text = bullet
    p.level = 1
# Add project image if available
image_path = Path("PROJECT IMAGE 1.webp")
if image_path.exists():
    try:
        slide.shapes.add_picture(str(image_path), Inches(6), Inches(1.5), width=Inches(3.5))
    except Exception:
        # convert to PNG if needed
        img = Image.open(image_path)
        png_path = Path("project_image_converted.png")
        img.save(png_path)
        slide.shapes.add_picture(str(png_path), Inches(6), Inches(1.5), width=Inches(3.5))

# Slide 7: Conclusion
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Conclusion"
body = slide.shapes.placeholders[1].text_frame
body.text = "THE SAFE SPACE HUB is a complete mental wellness prototype that blends secure access, assessment tools, and compassionate design."
body.add_paragraph().text = "It supports users with mood tracking, clinical-style assessments, emergency resources, and accessible UX."
body.add_paragraph().text = "This project is a strong foundation for further development into a full HIPAA-ready mental health platform."

presentation.save("THE_SAFE_SPACE_HUB_Presentation.pptx")
print("Presentation created: THE_SAFE_SPACE_HUB_Presentation.pptx")
